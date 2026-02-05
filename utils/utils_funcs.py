# -*- coding: utf-8 -*-
"""
General utility functions to analyze 129Xe ventilation images.
Used in the ventilation defect/hyper percentage (vdp/hvp) calculations.

Created on Wed Nov 29 15:12:06 2023
@author: Riaz Hussain, PhD
"""
import os
import warnings
import sys
import re
import csv
import nibabel as nib
import numpy as np
import matplotlib.pyplot as plt
from matplotlib.colors import ListedColormap, BoundaryNorm, Normalize
from matplotlib.cm import ScalarMappable
from scipy.ndimage import zoom, gaussian_filter
from pptx import Presentation
from pptx.util import Inches

class Tee:
    """Duplicates the output to multiple files in real-time."""
    def __init__(self, *files):
        """Initialize the Tee object with a list of files to duplicate output to."""
        self.files = files
    def write(self, obj):
        """Write the given object to all files in real-time."""
        for f in self.files:
            f.write(obj)
            f.flush()
    def flush(self):
        """Flush the output to all files in real-time."""
        for f in self.files:
            f.flush()

class ConsoleToFile:
    """Redirects console output to both the console and a log file."""
    def __init__(self, log_file_path):
        """Initialize the ConsoleToFile object with a log file path."""
        self.console = sys.stdout
        self.log_file = open(log_file_path, 'w', encoding="utf-8")
    def start(self):
        """Start redirecting console output to the log file."""
        sys.stdout = Tee(self.console, self.log_file)
    def stop(self):
        """Stop redirecting console output to the log file and close the log file."""
        sys.stdout = self.console
        self.log_file.close()

def reorient(array_in):
    """
    A function that reorients the array for montage.
    Args:
        array_in (ndarray): 3D array, in form N_pixel x N_pixel x N_slice.
    """
    output = np.flip(np.rot90(array_in), 0)
    return output

def process_nifti_file(input_type, patterns, subj_dir):
    """
    Search for a NIfTI file in the given directory that matches the provided filename patterns,
    and process the file by reorienting its data and return the data as output.
    Parameters:
    -----------
    input_type (str): Input type specifying pattern ('Non_corr', 'N4_corr', or 'FA_corr').
    patterns (dict): The dictionary which specifies the strings to be used.
    subj_dir (str): The path to the directory to search for the NIfTI file.
    Returns:
    --------
    data (numpy.ndarray): The reoriented data of the NIfTI file.
    Raises:
    -------
    FileNotFoundError
        If no NIfTI file is found in the directory that matches any of the provided patterns.
    ValueError
        If multiple NIfTI files are found in the directory that match the provided patterns.
    """
    # Compile patterns into regular expressions
    compiled_patterns = {pattern_type: [re.compile(pattern) for pattern in pattern_list]
                         for pattern_type, pattern_list in patterns.items()}
    if input_type not in compiled_patterns.keys():
        print("Input type not found in patterns. Returning None")
        return None
    filename_patterns = compiled_patterns[input_type]
    matching_files = []
    for filename in os.scandir(subj_dir):
        if filename.is_file():
            for pattern in filename_patterns:
                if pattern.match(filename.name):
                    matching_files.append(filename)
    if len(matching_files) == 0:
        raise FileNotFoundError(f"\nNothing matching the provided patterns in {subj_dir}")
    if len(matching_files) > 1:
        raise ValueError(f"\nMultiple files matching the provided names in {subj_dir}")
    print(f"\nLoading {matching_files[0].name}")
    nifti_load = nib.load(matching_files[0].path)
    nifti_data = nifti_load.get_fdata()
    reoriented_data = reorient(nifti_data)
    return reoriented_data

def med_filter(array_in, nome="array"):
    """
    Applies a median filter to a numpy array.
    Parameters:
    array_in (ndarray): Numpy array to be filtered.
    nome (str): optional name of the array to print out - italian of name = nome
    Returns:
    output_array (ndarray): Numpy array with the median filter applied.
    Notes:
    Median filter is applied to the input array with a window size of 3x3 in each dimension.
    If the array_in has NaN values, they are replaced with 0 in the output array.
    If the array_in = 2D, median filter is applied to each element of the array.
    If the array_in = 3D, median filter is applied to each 2D slice of the array.
    If the array_in = 4D, median filter is applied to each 2D slice of each 3D volume of array.
    """
    if array_in.ndim == 2:
        print(f"\nApplying median filter to 2D {nome}")
        output_array = np.zeros_like(array_in)
        for i in range(1, array_in.shape[0]-1):
            for j in range(1, array_in.shape[1]-1):
                #Window Size 3 x 3
                list_ = array_in[i-1:i+2, j-1:j+2]
                output_array[i,j] = np.nanmedian(list_)
        output_array[np.isnan(output_array)] = 0

    if array_in.ndim == 3:
        print(f"\nApplying median filter to 3D {nome}")
        output_array = np.zeros_like(array_in)
        rows, cols, depth = array_in.shape
        for k in range(depth):
            for i in range(1, rows - 1):
                for j in range(1, cols - 1):
                    list_ = array_in[i-1:i+2, j-1:j+2, k]
                    output_array[i, j, k] = np.nanmedian(list_)
        output_array[np.isnan(output_array)] = 0

    if array_in.ndim == 4:
        print(f"\nApplying median filter to 4D {nome}")
        output_array = np.zeros((array_in.shape[0], array_in.shape[1], array_in.shape[2],
                                 array_in.shape[3]))
        for b in range(array_in.shape[3]):
            for k in range(array_in.shape[2]):
                for i in range(1, array_in.shape[0]-1):
                    for j in range(1, array_in.shape[1]-1):
                        list_ = array_in[i-1:i+2, j-1:j+2, k, b]
                        output_array[i, j, k, b] = np.nanmedian(list_)
        output_array[np.isnan(output_array)] = 0

    return output_array

def create_srm_array(image_in, mask_file=None, param='mean'):
    """
    Calculate the relative signal with respect to the median of a 3D image.
    Parameters:
    - image_in: input numpy array of 3D MR image
    - Optional: binary mask image, parameter (mean or median), path to save file,
                and type of data being analyzed
    Returns:
    - rel_sig_img: numpy array, relative mean/median-to-signal ratio image
    """
    param_value = np.median(image_in) if param == 'median' else np.mean(image_in)
    rel_sig_img = image_in / param_value
    if mask_file is not None and mask_file.any():
        rel_sig_img *= mask_file
        slices = np.flatnonzero(np.sum(mask_file, axis=(0, 1)))
    else:
        slices = list(range(rel_sig_img.shape[2]))
    rel_sig_img[rel_sig_img == 0] = np.nan
    montage_array = np.hstack([rel_sig_img[:, :, i] for i in slices])
    # montage_array = (montage_array - np.nanmin(montage_array)
    #                  ) / (np.nanmax(montage_array) - np.nanmin(montage_array))
    print(f"{param}-to-signal ratio array created.")
    return montage_array

def calc_slicewise_srm(msked_norm_image, msr_dir, subj, correc, param='mean'):
    """calculate and save slice-wise SRM (signal-relative-to-mean)"""
    param_value = (np.nanmedian(msked_norm_image) if param == 'median'
                   else np.nanmean(msked_norm_image))
    slicewise_norm_srm = []
    for a_slice in range(msked_norm_image.shape[2]):
        try:
            with warnings.catch_warnings():
                warnings.simplefilter("ignore", category=RuntimeWarning)
                srm_r = msked_norm_image[:,:, a_slice] / param_value
                slice_srm_r = np.nanmedian(srm_r) if param == 'median' else np.nanmean(srm_r)
            slicewise_norm_srm.append(slice_srm_r)
        except ValueError as err:
            print(f"Error calculating normalized SRM for slice {a_slice}: {err}")
            continue
    slicewise_norm_srm = np.array(slicewise_norm_srm)
    srm_no_nan = slicewise_norm_srm[~np.isnan(slicewise_norm_srm)]
    # srm_no_nan /= np.max(srm_no_nan)
    save_slicewise_srm(srm_no_nan, msr_dir, subj, correc)

def save_slicewise_srm(slicewise_srm_vals, msr_dir, subjid, corr):
    """Save slice-wise SNR"""
    slwise_fname = os.path.join(msr_dir, f"{subjid}_{corr}_slicewiseSRM.csv")
    distance = np.arange(1.5, 1.5 * (len(slicewise_srm_vals) + 1), 1.5)
    ##Save the srm array with distance
    with open(slwise_fname, mode='w', newline='', encoding="utf-8") as output_file:
        writer = csv.writer(output_file)
        writer.writerow(["distance", "srm"])
        ##Use zip to pair corresponding values from distance and slicewise_srm_vals
        for val1, val2 in zip(distance, slicewise_srm_vals):
            writer.writerow([val1, val2])

def create_montage_array(montage_in, slices=None, msk_file=None):
    """
    Creates a 2D N_pixel x (N_slice * N_pixel (x N_Channels)) array from a
    3D (4D) N_pixel x N_pixel x N_slice (x N_channel) array.
    Args:
        montage_in (ndarray): 3D/4D array, in form N_pixel x N_pixel x N_slice (x N_channel).
        msk_file (ndarray): 3D array with the same (3) dimensions as montage_in (optional).
        slices (list of ints): Slices to make into montage. Defaults to middle 7 slices.
        Other options: 'all' plots all the slices if msk_file=None. 
                        'all' plots non-zero(=masked) slices if mask provided.
    Returns:
        montage_output (ndarray): Numpy array for making montage.
    """
    if slices is None and msk_file is None:
        mid_im = montage_in.shape[2] // 2
        slices = [mid_im - 3, mid_im - 2, mid_im - 1, mid_im, mid_im + 1, mid_im + 2, mid_im + 3]
    elif slices == "all" and msk_file is None:
        slices = list(range(montage_in.shape[2]))
    else:
        slices = np.flatnonzero(np.sum(msk_file, axis=(0, 1)))
    if np.ndim(montage_in) == 3:
        montage_output = montage_in[:, :, slices[0]]
        for i in slices[1:]:
            montage_output = np.hstack((montage_output, montage_in[:, :, i]))
    elif np.ndim(montage_in) == 4:
        montage_output = montage_in[:, :, slices[0], :]
        for i in slices[1:]:
            montage_output = np.hstack((montage_output, montage_in[:, :, i, :]))
    return montage_output

def montage_plot_3d(montage_in, prtn=None, clr=False, cmapp='jet'): # msk_file=None,
    """To plot the montage image
        montage_in (ndarray): 3D/4D array, in form N_pixel x N_pixel x N_slice (x N_channel).
        msk_file (ndarray): (optional) 3D binary mask array used for proton image matching.
        prtn (bool) = if provided, plots the proton image as background of other montage
        clr (bool) = if True, plots a jet map of 3d montage array
    """
    fig, ax = plt.subplots(figsize=(16, 3), dpi=300)
    if prtn is not None and np.any(prtn):
        prtn_montage = np.copy(prtn)
        blurred_proton = gaussian_filter(prtn_montage, sigma=6)
        blurred_proton /= np.nanmax(blurred_proton)
        ax.imshow(blurred_proton, cmap='gray', alpha=1)
        if clr:
            cmp = plt.get_cmap(cmapp)
            img = ax.imshow(montage_in, cmap=cmp, alpha=1)
            cax = fig.add_axes([ax.get_position().xmax + 0.002, ax.get_position().y0,
                                0.01, ax.get_position().height])
            nrm = Normalize(vmin=0, vmax=1)
            cbar = plt.colorbar(img, cax=cax, ticks=[0, 0.5, 1], norm=nrm)
            cbar.ax.set_yticklabels(labels=cbar.ax.get_yticklabels(),weight='bold',
                                    fontsize=12)
            ax.set_xticks([])
            ax.set_yticks([])
            ax.tick_params(axis='both', which='both', length=0)
        else:
            ax.imshow(montage_in, cmap='gray', alpha=1)
            ax.tick_params(axis='both', which='both', length=0)
            ax.set_xticks([])
            ax.set_yticks([])
    else:
        if clr:
            cmp = plt.get_cmap(cmapp)
            img = ax.imshow(montage_in, cmap=cmp, alpha=1)
            cax = fig.add_axes([ax.get_position().xmax + 0.002, ax.get_position().y0,
                                0.01, ax.get_position().height])
            nrm = Normalize(vmin=0, vmax=1)
            cbar = plt.colorbar(img, cax=cax, ticks=[0, 0.5, 1], norm=nrm)
            cbar.ax.set_yticklabels(labels=cbar.ax.get_yticklabels(),weight='bold',
                                    fontsize=12)
            cmp.set_bad(color="black")
            ax.set_xticks([])
            ax.set_yticks([])
            ax.tick_params(axis='both', which='both', length=0)
        else:
            cmp = plt.get_cmap('gray')
            ax.imshow(montage_in, cmap=cmp, alpha=1)
            cmp.set_bad(color="black")
            ax.tick_params(axis='both', which='both', length=0)
            ax.set_xticks([])
            ax.set_yticks([])
    plt.show(block=False)
    return fig

def montage_plot_4d(montage_in):
    """To plot the 4D (3D image + color overlaid) montage image
        montage_in (ndarray): 4D array, in form N_pixel x N_pixel x N_slice (x N_channel),
        should be processed through create_montage_array() function already
    """
    fig, ax = plt.subplots(figsize=(16, 3), dpi=300)
    ax.imshow(montage_in, alpha=1)
    plt.axis('off')
    return fig

def save_image(fig, save_path, data_type=None, im_type=None):
    """To save image in the provided path
        save_path (str): path to directory for saving the image.
        data_type (str): type of data (None (default), Non_corr, N4_corr, FA_corr)
        im_type (str): optional image data type (e.g. raw, mask, defect, montage)
    """
    if data_type is not None:
        filename = os.path.join(save_path, f"{data_type}_img_{im_type}.png")
    else:
        ## benaam from Urdu means without/unknown name
        filename = os.path.join(save_path, 'benaam_img.png')
    i = 0
    while os.path.exists(filename):
        i += 1
        filename = f"{os.path.splitext(filename)[0]}_{i}.png"
    fig.savefig(filename, dpi=300, bbox_inches='tight', pad_inches=0)
    plt.show(block=False)

def match_img_dimensions(vent_im, proton_image):
    """
    Match the dimensions (width, height, and depth) of vent and proton images,
    while centering them, with zero padding or downsampling of proton image.
    Args:
        vent_im: ventilation image (NumPy array).
        proton_image: proton image (NumPy array).
    Returns:
        matched_proton_image: proton image after resizing,
                    to match the dimensions of vent image.
    """
    # Get the dimensions of the input images
    vent_shape = np.array(vent_im.shape)
    proton_shape = np.array(proton_image.shape)
    # Create new arrays with the target shape, filled with zeros
    matched_proton_image = np.zeros(vent_shape)
    # Copy the input images into the new arrays with zero padding or downsampling
    if np.any(vent_shape < proton_shape):
        # Downsample the proton image to match dimensions of the vent image
        zoom_factors = vent_shape / proton_shape
        matched_proton_image = zoom(proton_image, zoom_factors, order=1)
    else:
        # Calculate the starting indices for centering the  ptoton image
        start_proton = (vent_shape - proton_shape) // 2
        # Use zero padding to match the dimensions of the proton image
        matched_proton_image[start_proton[0]:start_proton[0] + proton_shape[0],
                             start_proton[1]:start_proton[1] + proton_shape[1],
                             start_proton[2]:start_proton[2] + proton_shape[2]] = proton_image
    return matched_proton_image

def create_new_folder(initial_dir, folder_name='results', not_create_if_exist=False):
    """
    Create a new folder with a new (unique name) in the specified directory.
    If folder_name already exists and is not empty, it'll append a number to it.
    Args:
        initial_dir (str): The initial directory where the new folder will be created.
        folder_name (str): The desired name of the folder. Default: 'results'
    Returns:
        str: The path to the newly created folder or the old folder (if not empty).
    """
    counter = 1
    new_folder = os.path.join(initial_dir, folder_name)
    while os.path.exists(new_folder):
        if os.listdir(new_folder) and not_create_if_exist is False:
            new_folder = os.path.join(initial_dir, folder_name + "_" + str(counter))
            counter += 1
            print("Old folder is not empty. Will append a number")
        elif os.listdir(new_folder) and not_create_if_exist is True:
            print("Folder already exists. Not creating new!")
            return new_folder
        else:
            print("Folder already exists and is empty...")
            return new_folder
    os.makedirs(new_folder)
    print(f"New folder created: {new_folder}")
    return new_folder

def save_txt_file(data, data_path, file_name='data.txt', txt_heading=None,
              subj_name=None):
    """
    Save data (must be as string) as a text file to the specified directory.
    If the file already exists, it'll create a new file with a number appended to the filename.
    Args:
        data (str): Content to be saved as a text file.
        data_path (str): Directory where the text file will be saved.
        file_name (str): Name of the text file. Default: 'file.txt'
        txt_heading (str): if provided inserts the heading on top of data
        subj_name (str): if provided, adds a column with subject name before data
    Returns:
        str: The path to the saved text file.
    """
    file_path = os.path.join(data_path, file_name)
    counter = 1
    while os.path.exists(file_path):
        file_name_parts = os.path.splitext(file_name)
        file_name = f"{file_name_parts[0]}_{counter}{file_name_parts[1]}"
        file_path = os.path.join(data_path, file_name)
        counter += 1
        print("file already exists. Will append a number")
    with open(file_path, 'w', encoding="utf-8") as file:
        if txt_heading:
            file.write(f"#\t{txt_heading}\n")
        if subj_name:
            file.write(f"{subj_name}\t")
        file.write(str(data))
    print(f"\nText file saved: {file_path}")

def append_save_txt2csv(dir_path, f_name, csv_heading, correction=None):
    """To append data txt files inside a csv file
    Args:
        main_dir_path (string): main directory containing folders of text files
        f_name (string): file name to be appended
        csv_heading (tuple): Tuple of heading strings for the CSV file
        correction: name of bias correction applied to MR images
    """
    # Initialize an empty list to store the data
    data = []
    im_type, analysis_type = None, None
    # Traverse the main directory and find files containing the given name
    for root, _, files in os.walk(dir_path): # _ = dirs
        if f_name in files:
            im_type, analysis_type = (correction or f_name.split("_")[0], f_name.split("_")[-2])
            text_file_path = os.path.join(root, f_name)
            print(f"{text_file_path}\nimage type {im_type}, analysis type {analysis_type}")
            try:
                # Read the data from the txt file and append it to the list
                with open(text_file_path, 'r', encoding="utf-8") as file:
                    for line in file.readlines():
                        if line.startswith('#'):  # skip header row
                            continue
                        subject_id, data_values = line.split('\t')
                        print(subject_id)
                        data.append([subject_id] + [float(value) for value in
                                                    data_values.strip('[]\n').split(',')])
            except (FileNotFoundError, ValueError) as e:
                print(f"Error processing file {text_file_path}: {e}")
                continue  # Skip to the next iteration if an error occurs
    ##Save csv file with fetched data/heading
    if im_type is not None and analysis_type is not None:
        save_csv_file(data, dir_path, im_type, analysis_type, csv_heading)

def save_csv_file(data_in, main_dir, im_type, analysis_type, heading):
    """
    Create and save a CSV file
    heading (string): header text
    """
    csv_path = os.path.join(main_dir, f"{im_type}_{analysis_type}_analysis_results.csv")
    with open(csv_path, 'w', newline='', encoding="utf-8") as file:
        writer = csv.writer(file)
        writer.writerow(heading)
        for row in data_in:
            writer.writerow(row)
    print(f"csv file saved: {csv_path}")

def append_save_snr_txts(main_dir, heading, correction, analysis_type):
    """to append data txt files inside a csv file
    Args:
        main_dir (string): main directory containing folders of text files
        heading (string): header text
        correction: name of bias correction applied to MR images
        analysis_type: what the data analysed for/with
    """
    # Initialize an empty list to store the data
    data = []
    txt_file_name = f'{correction}_{analysis_type}.txt'
    # Traverse the main directory and find files containing the given name
    for root, _, files in os.walk(main_dir):
        if txt_file_name in files:
            txt_file_path = os.path.join(root, txt_file_name)
            print(txt_file_path)
            print(f"image type {correction}, analysis type {analysis_type}")
            # #Read the data from the txt file and append it to the list
            with open(txt_file_path, 'r', encoding="utf-8") as file:
                lines = file.readlines()
                for line in lines:
                    if line.startswith('#'):  # skip header row
                        continue
                    subject_id, data_values = line.split('\t')
                    print(subject_id)
                    data_values = data_values.strip('[]\n').split(', ')
                    data_values = [float(value) for value in data_values]
                    data.append([subject_id] + data_values)
    save_csv_file(data, main_dir, correction, analysis_type, heading)

def fixed_cbar(num_colors):
    """Generate a colorbar with fixed colors and automatically calculated bounds"""
    if num_colors == 4:
        colors = ['#ff0000', '#ffb600', '#00ff00', '#0000ff']
    elif num_colors == 5:
        colors = ['#ff0000', '#ffb600', '#66b366', '#00ff00', '#0000ff']
    elif num_colors == 6:
        colors = ['#ff0000', '#ffb600', '#66b366', '#00ff00', '#0091b5', '#0000ff']
    else:
        raise ValueError("Only 4, 5 or 6 bin colors are supported so far.")
    # Create bounds based on the number of colors
    bounds = list(range(num_colors + 1))
    # Create colormap
    cbar_colors = ListedColormap(colors)
    s_map = ScalarMappable(cmap=cbar_colors, norm=BoundaryNorm(bounds, cbar_colors.N))
    return s_map

def binning_cmap(bin_percents, num_colors, prtn_arr=None):
    """Create a colormap for binned glb images based on the number of colors"""
    if num_colors == 4:
        colors = ['#00000000','#ff0000',"#ffb600","#00ff00", "#0000ff"]
        vent_counts = [bin_percents[0], bin_percents[1], bin_percents[2],
                       bin_percents[3]]
    elif num_colors == 5:
        colors = ['#00000000','#ff0000',"#ffb600","#66b366","#00ff00","#0000ff"]
        vent_counts = [bin_percents[0], bin_percents[1], bin_percents[2],
                       bin_percents[3], bin_percents[4]]
    elif num_colors == 6:
        colors = ['#00000000','#ff0000',"#ffb600","#66b366","#00ff00","#0091b5","#0000ff"]
        vent_counts = [bin_percents[0], bin_percents[1], bin_percents[2],
                       bin_percents[3], bin_percents[4], bin_percents[5]]
    else:
        raise ValueError("Only 4, 5, or 6 num_colors in colormap are supported")
    if prtn_arr is not None and np.any(prtn_arr):
        colors[0] = '#00000000'
    else:
        colors[0] = '#000000'
    my_cmap = ListedColormap(colors)
    color_bounds = [0]
    for i, count in enumerate(vent_counts, start=1):
        if count > 0:
            color_bounds.append(i)
    ## Ensure indices in color_bounds are within the range of my_cmap.colors
    max_index = len(my_cmap.colors) - 1
    color_bounds = [i for i in color_bounds if i <= max_index]
    ## Adjust colormap to remove colors for values with a count of 0
    my_cmap = ListedColormap([my_cmap.colors[i] for i in color_bounds])
    return my_cmap

def create_ppt(parent_dir, image_names, image_titles):
    """To create a ppt and place images on it.
    """
    # Create a new PowerPoint presentation
    prs = Presentation()
    # Set the slide size to 16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    # Add a slide with a blank layout
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Use index 6 for a blank slide layout
    # Add images to the slide vertically stacked
    top = Inches(0.5)
    image_dict = dict(zip(image_names, image_titles))  #Dictionary to map image names to titles
    for dirpath, _, filenames in os.walk(parent_dir):
        # Look for PNG images with filenames that contain any of the image_names
        image_files = [f for f in filenames if f.endswith('.png')
                       and any(name in f for name in image_names)]
        for image_file in image_files:
            image_path = os.path.join(dirpath, image_file)
            title = image_dict[image_file]  # Get the title from the dictionary
            slide.shapes.add_picture(image_path, Inches(0), top, width=prs.slide_width, height=None)
            title_shape = slide.shapes.add_textbox(Inches(0), top -
                                                   Inches(0.4), Inches(2), Inches(0.4))
            title_shape.text_frame.text = title
            top += Inches(1.75)
    # Save the PowerPoint presentation in the parent directory
    pptx_path = os.path.join(parent_dir, 'result_images_ppt.pptx')
    prs.save(pptx_path)
    print(f'Saved PowerPoint: {pptx_path}')

#%%
